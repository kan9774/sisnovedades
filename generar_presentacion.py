#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Genera la presentación SIS-Novedades para usuarios finales.
Stack: python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paleta corporativa (azul militar + dorado) ──
COLOR_PRIMARIO   = RGBColor(0x0B, 0x25, 0x45)  # azul oscuro
COLOR_SECUNDARIO = RGBColor(0xFF, 0xD2, 0x00)  # dorado
COLOR_TEXTO      = RGBColor(0x1A, 0x1A, 0x1A)  # casi negro
COLOR_GRIS       = RGBColor(0x66, 0x66, 0x66)  # gris medio
COLOR_FONDO      = RGBColor(0xFF, 0xFF, 0xFF)  # blanco
COLOR_CLARO      = RGBColor(0xF0, 0xF4, 0xF8)  # azul muy claro
COLOR_VERDE      = RGBColor(0x1E, 0x7E, 0x34)  # verde
COLOR_ROJO       = RGBColor(0xC0, 0x39, 0x2B)  # rojo

# ── Fuentes ──
FONT_TITULO   = 'Calibri'
FONTCUERPO    = 'Calibri'
FONT_NOTAS    = 'Calibri'

def configurar_slide(prs, layout_idx=0):
    """Crea un slide con el layout indicado."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def agregar_titulo(slide, texto, ancho=Inches(10), alto=Inches(1.2),
                   izquierda=Inches(0.5), top=Inches(0.3),
                   color=COLOR_PRIMARIO, tamano=36, negrita=True):
    """Agrega un título con formato profesional."""
    txBox = slide.shapes.add_textbox(izquierda, top, ancho, alto)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.name = FONT_TITULO
    p.font.size = Pt(tamano)
    p.font.bold = negrita
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return txBox

def agregar_subtitulo(slide, texto, izquierda=Inches(0.5), top=Inches(1.5),
                      ancho=Inches(10), alto=Inches(0.6),
                      color=COLOR_GRIS, tamano=18):
    """Agrega un subtítulo."""
    txBox = slide.shapes.add_textbox(izquierda, top, ancho, alto)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.name = FONTCUERPO
    p.font.size = Pt(tamano)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return txBox

def agregar_linea_decorativa(slide, izquierda=Inches(0.5), top=Inches(1.55),
                             ancho=Inches(1.5), alto=Emu(0), color=COLOR_SECUNDARIO):
    """Agrega una línea decorativa dorada bajo el título."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, izquierda, top, ancho, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def agregar_bullets(slide, items, izquierda=Inches(0.8), top=Inches(2.2),
                    ancho=Inches(9), alto=Inches(5), tamano=18,
                    color=COLOR_TEXTO, espaciado=Pt(8)):
    """Agrega viñetas con formato profesional."""
    txBox = slide.shapes.add_textbox(izquierda, top, ancho, alto)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.font.name = FONTCUERPO
        p.font.size = Pt(tamano)
        p.font.color.rgb = color
        p.space_after = espaciado
        p.level = 0
        # Viñeta
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None:
            pPr.remove(buNone)
        # Agregar viñeta con carácter
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', FONT_TITULO)
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', '%02X%02X%02X' % color)

    return txBox

def agregar_notas(slide, texto):
    """Agrega notas del orador al slide actual."""
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = texto

def agregar_fecha(slide, texto, izquierda=Inches(0.5), top=Inches(7.2),
                  ancho=Inches(4), alto=Inches(0.4), tamano=12):
    """Agrega fecha en el pie."""
    txBox = slide.shapes.add_textbox(izquierda, top, ancho, alto)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = texto
    p.font.name = FONTCUERPO
    p.font.size = Pt(tamano)
    p.font.color.rgb = COLOR_GRIS
    p.alignment = PP_ALIGN.LEFT

def agregar_icono(slide, emoji, izquierda=Inches(0.5), top=Inches(0.5),
                  ancho=Inches(0.8), alto=Inches(0.8)):
    """Agrega un emoji como ícono visual."""
    txBox = slide.shapes.add_textbox(izquierda, top, ancho, alto)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.name = 'Segoe UI Emoji'
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER
    return txBox

# ═══════════════════════════════════════════════════════════
#  CREAR PRESENTACIÓN
# ═══════════════════════════════════════════════════════════

prs = Presentation()
# Configurar tamaño slide 16:9 (widescreen)
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ──────────────────────────────────────────────────────────
# SLIDE 1 — PORTADA
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 6)  # layout título

# Fondo azul
from pptx.oxml.ns import qn
from lxml import etree
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = COLOR_PRIMARIO

# Título principal
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'SIS-Novedades'
p.font.name = FONT_TITULO
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = COLOR_SECUNDARIO
p.alignment = PP_ALIGN.LEFT

# Línea dorada
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5), Inches(3), Pt(6))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_SECUNDARIO
shape.line.fill.background()

# Subtítulo
txBox = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'Sistema de Gestión de Novedades'
p.font.name = FONTCUERPO
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
p.alignment = PP_ALIGN.LEFT

p2 = tf.add_paragraph()
p2.text = 'Unidad BCOM1'
p2.font.name = FONTCUERPO
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p2.alignment = PP_ALIGN.LEFT

# Fecha
txBox = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(6), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = 'Presentación para personal administrativo'
p.font.name = FONTCUERPO
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

# ──────────────────────────────────────────────────────────
# SLIDE 2 — ¿QUÉ ES SIS-NOVEDADES?
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)  # layout título solo

agregar_titulo(slide, '¿Qué es SIS-Novedades?')
agregar_linea_decorativa(slide)
agregar_icono(slide, '💻', izquierda=Inches(11.5), top=Inches(0.5))

agregar_bullets(slide, [
    'Sistema digital para gestionar todas las novedades de la unidad',
    'Reemplaza las planillas manuales y los papeles dispersos',
    'Todo queda registrado, ordenado y accesible desde la computadora',
    'El corazón del sistema: la gestión de novedades del día',
], tamano=22)

agregar_notas(slide,
    'En términos simples, esto es donde cargamos y gestionamos todo lo que pasa en la unidad '
    'durante el día: los informes que llegan, los que salimos, las situaciones especiales. '
    'Antes todo estaba en papel o en planillas sueltas. Ahora todo queda en un solo lugar, '
    'ordenado por fecha y hora, y siempre disponible.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 3 — ANTES VS AHORA
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, '¿Qué problema resuelve?')
agregar_linea_decorativa(slide)
agregar_icono(slide, '⚖️', izquierda=Inches(11.5), top=Inches(0.5))

# Columna "Antes" (rojo)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'ANTES'
p.font.name = FONT_TITULO
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLOR_ROJO
p.alignment = PP_ALIGN.LEFT

items_ant = [
    'Novedades registradas en papel, fácil de perder o desordenar',
    'Correos enviados uno por uno, manual y propenso a errores',
    'Nadie sabía si los correos llegaron a destino',
]
for i, item in enumerate(items_ant):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = '✗  ' + item
    p.font.name = FONTCUERPO
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_ROJO
    p.space_after = Pt(12)

# Columna "Ahora" (verde)
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'AHORA'
p.font.name = FONT_TITULO
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLOR_VERDE
p.alignment = PP_ALIGN.LEFT

items_ahora = [
    'Todo se carga en pantalla, queda registrado automáticamente',
    'El sistema arma el reporte diario y lo envía solo por correo',
    'Control automático de correos no entregados',
]
for i, item in enumerate(items_ahora):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = '✓  ' + item
    p.font.name = FONTCUERPO
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_VERDE
    p.space_after = Pt(12)

agregar_notas(slide,
    'Imaginen esto: fin de turno, tenés que armar el reporte del día con todo lo que se registró. '
    'Antes armaban a mano, lo pasaban a PDF, y luego tenían que enviarlo uno por uno a unas 30 personas. '
    'Si se les olvidaba alguien, o si un correo rebotaba y no se enteraban… bueno, ahí empezaba el problema. '
    'Ahora el sistema hace todo eso automáticamente.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 4 — CARGA DE NOVEDADES
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, 'Carga de novedades del día')
agregar_linea_decorativa(slide)
agregar_icono(slide, '📝', izquierda=Inches(11.5), top=Inches(0.5))

agregar_bullets(slide, [
    'Se registra todo lo que llega y sale durante la guardia',
    'Cada novedad tiene hora, número, asunto y clasificación',
    'Se clasifican como: Rutinario, Prioritario, Urgente o Destello',
    'Se organizan automáticamente por tipo y dirección (Recibido/Expedido)',
    'Se pueden editar o corregir mientras la guardia está abierta',
], tamano=22)

agregar_notas(slide,
    'Durante la guardia, el escribiente va cargando cada novedad: un radio que llegó, '
    'un correo que se envió, una orden que se recibió. Cada una lleva hora, número y asunto. '
    'El sistema las ordena solas por tipo y por si son recibidas o enviadas. '
    'Y si se equivocaron, las pueden corregir mientras la guardia esté abierta.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 5 — REPORTE DIARIO
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, 'Generación del reporte diario')
agregar_linea_decorativa(slide)
agregar_icono(slide, '📄', izquierda=Inches(11.5), top=Inches(0.5))

agregar_bullets(slide, [
    'Con todas las novedades cargadas, se arma el reporte del día',
    'El sistema genera un PDF con todo el tráfico registrado',
    'Se puede revisar antes de enviar, con vista previa',
    'Incluye datos de la guardia: fecha, comandante, oficial de día',
], tamano=22)

agregar_notas(slide,
    'Cuando llega el momento de enviar el reporte, el sistema toma todas las novedades '
    'que se cargaron durante la guardia y las junta en un solo documento PDF bien ordenado. '
    'Se puede revisar antes de enviar para asegurarse de que todo esté correcto. '
    'El reporte incluye los datos de la guardia: quién fue el comandante, el oficial de día, la fecha.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 6 — ENVÍO AUTOMÁTICO
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, 'Envío automático por correo')
agregar_linea_decorativa(slide)
agregar_icono(slide, '📧', izquierda=Inches(11.5), top=Inches(0.5))

agregar_bullets(slide, [
    'Se eligen los destinatarios (por lista o por grupo predefinido)',
    'El sistema envía el PDF a todos a la vez, sin intervención manual',
    'Aproximadamente 30 destinatarios por envío',
    'Se puede enviar con adjuntos embebidos o como archivo ZIP',
    'Confirmación instantánea de que el envío se procesó',
], tamano=22)

agregar_notas(slide,
    'Acá está la gran ventaja: uno elige a quién se lo envía — puede ser una lista personalizada '
    'o un grupo predefinido como "todas las oficinas"—, y el sistema se encarga de mandar el PDF '
    'a las 30 o más personas. Sin tener que copiar y pegar direcciones, sin tener que adjuntar '
    'el archivo 30 veces. Uno hace clic y listo.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 7 — CORREOS NO ENTREGADOS
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, 'Control de correos no entregados')
agregar_linea_decorativa(slide)
agregar_icono(slide, '⚠️', izquierda=Inches(11.5), top=Inches(0.5))

agregar_bullets(slide, [
    'Si un correo rebota (no llega a destino), el sistema lo detecta solo',
    'Queda registrado en una lista de correos fallidos con el motivo',
    'Se puede reenviar manualmente desde el sistema con un clic',
    'Nadie tiene que estar revisando manualmente si los correos llegaron',
], tamano=22)

agregar_notas(slide,
    'Y acá hay otra cosa que antes era un dolor de cabeza: si un correo rebotaba, '
    'nadie se enteraba hasta que el destinatario se quejaba. Ahora, si un correo no llega, '
    'el sistema lo detecta automáticamente y lo registra con el motivo del fallo. '
    'Desde ahí se puede reenviar con un clic. Y hay un aviso visible que muestra cuántos '
    'correos fallidos hay, para que nadie se quede con la duda.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 8 — TRAZABILIDAD
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, 'Toda la información queda registrada')
agregar_linea_decorativa(slide)
agregar_icono(slide, '🗂️', izquierda=Inches(11.5), top=Inches(0.5))

agregar_bullets(slide, [
    'Cada novedad queda guardada con fecha y hora exacta',
    'Se puede consultar el historial de cualquier día anterior',
    'Se puede ordenar por fecha, hora, tipo o clasificación',
    'El sistema registra quién cargó cada novedad y cuándo',
    'Las novedades se pueden recuperar desde la papelera',
], tamano=22)

agregar_notas(slide,
    'Todo lo que se carga queda guardado para siempre. Si mañana necesitan saber qué se registró '
    'el 15 de marzo del año pasado, lo pueden buscar sin problema. El sistema también registra '
    'quién hizo cada cosa: quién cargó la novedad, quién la corrigió, quién envió el correo. '
    'Y si por error se borra algo, está en la papelera y se puede recuperar.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 9 — BENEFICIOS
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, '¿Qué gano yo con esto?')
agregar_linea_decorativa(slide)
agregar_icono(slide, '✅', izquierda=Inches(11.5), top=Inches(0.5))

agregar_bullets(slide, [
    '✅  Menos tiempo: el reporte se arma y envía solo',
    '✅  Menos errores: nada se pierde, nada se olvida',
    '✅  Información siempre ordenada y accesible',
    '✅  Seguimiento de correos sin tener que estar revisando',
    '✅  Historial completo para consultas futuras',
], tamano=22, color=COLOR_TEXTO)

agregar_notas(slide,
    'En resumen: esto les va a quitar trabajo manual. Ya no tienen que armar reportes a mano '
    'ni enviar correos uno por uno. La información queda ordenada y siempre disponible. '
    'Si tienen que buscar algo de meses atrás, lo encuentran en segundos. '
    'Y el sistema les avisa si algo no se envió bien, para que puedan corregirlo.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 10 — REPASO OTROS MÓDULOS
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, 'Otros módulos del sistema')
agregar_linea_decorativa(slide)
agregar_icono(slide, '🧩', izquierda=Inches(11.5), top=Inches(0.5))

# Tabla de módulos
from pptx.util import Inches, Pt
table_data = [
    ('Módulo', 'Qué hace'),
    ('🚗 Vehículos y Conductores', 'Registro de vehículos, habilitación de conductores y actas'),
    ('🚛 Salidas de Vehículos', 'Control de cuándo sale y vuelve cada vehículo'),
    ('🏢 Oficinas y Organismos', 'Estructura organizativa de la unidad'),
    ('🔐 Roles y Permisos', 'Cada usuario ve solo lo que le corresponde'),
    ('💾 Respaldo Automático', 'Copias de seguridad periódicas, información protegida'),
]

rows, cols = len(table_data), 2
tbl = slide.shapes.add_table(rows, cols, Inches(1), Inches(2.2), Inches(11), Inches(4.5)).table

# Anchos de columnas
tbl.columns[0].width = Inches(4)
tbl.columns[1].width = Inches(7)

for i, row_data in enumerate(table_data):
    for j, cell_text in enumerate(row_data):
        cell = tbl.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONTCUERPO
            if i == 0:
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLOR_FONDO
            else:
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = COLOR_TEXTO
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# Estilos de tabla
from pptx.oxml.ns import qn
from lxml import etree
for i, row in enumerate(tbl.rows):
    for j, cell in enumerate(row.cells):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        if i == 0:
            # Header
            bg_color = COLOR_PRIMARIO
        else:
            bg_color = COLOR_FONDO if i % 2 == 1 else COLOR_CLARO
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        if i == 0:
            srgbClr.set('val', '%02X%02X%02X' % bg_color)
        else:
            srgbClr.set('val', '%02X%02X%02X' % bg_color)

agregar_notas(slide,
    'Además de las novedades, el sistema tiene otros módulos que complementan la gestión. '
    'Está el control de vehículos y conductores, las salidas de vehículos, la estructura de oficinas, '
    'los permisos de cada usuario, y respaldos automáticos para que la información nunca se pierda. '
    'Pero el día a día, lo que más van a usar es el módulo de novedades.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 11 — PRÓXIMOS PASOS
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 1)

agregar_titulo(slide, 'Cómo empezar a usarlo')
agregar_linea_decorativa(slide)
agregar_icono(slide, '🚦', izquierda=Inches(11.5), top=Inches(0.5))

agregar_bullets(slide, [
    'El sistema está disponible en: [URL o dirección de acceso]',
    'Se accede con el usuario y contraseña que ya tienen',
    'El soporte técnico está en: [contacto de soporte]',
    'Se recomienda hacer una prueba con una guardia de ejemplo',
    'Cualquier duda, consultar con [responsable del proyecto]',
], tamano=22)

agregar_notas(slide,
    'Para empezar a usarlo, solo necesitan entrar a la dirección del sistema con su usuario '
    'y contraseña habitual. Les recomiendo que prueben cargando una guardia de ejemplo para '
    'familiarizarse con la interfaz. Si tienen alguna duda o problema, el contacto de soporte '
    'está [indicar contacto]. Estamos acá para acompañarlos en la transición.'
)

# ──────────────────────────────────────────────────────────
# SLIDE 12 — CIERRE
# ──────────────────────────────────────────────────────────
slide = configurar_slide(prs, 6)  # layout título

# Fondo azul
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = COLOR_PRIMARIO

# Título
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '¿Preguntas?'
p.font.name = FONT_TITULO
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = COLOR_SECUNDARIO
p.alignment = PP_ALIGN.LEFT

# Subtexto
txBox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'Gracias por su atención'
p.font.name = FONTCUERPO
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
p.alignment = PP_ALIGN.LEFT

p2 = tf.add_paragraph()
p2.text = 'El sistema está listo para usar'
p2.font.name = FONTCUERPO
p2.font.size = Pt(22)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p2.alignment = PP_ALIGN.LEFT

# ──────────────────────────────────────────────────────────
# GUARDAR
# ──────────────────────────────────────────────────────────
output_path = 'presentacion_sisnovedades.pptx'
prs.save(output_path)
print('Presentacion guardada: ' + output_path)
print('   ' + str(len(prs.slides)) + ' diapositivas generadas')
